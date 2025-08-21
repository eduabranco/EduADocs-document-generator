"""
Document Generators Module
Handles generation of different document types (exercises, PowerPoint, summaries).
"""

from llm_config import get_llm_response
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import io


def generate_document(params):
    """Main document generation coordinator"""
    try:
        doc_type = params["doc_type"]
        
        if doc_type == "Exercise List":
            return generate_exercises(params)
        elif doc_type == "PowerPoint Presentation":
            return generate_powerpoint(params)
        elif doc_type == "Summary":
            return generate_summary(params)
        else:
            return {"success": False, "error": "Unknown document type"}
            
    except Exception as e:
        return {"success": False, "error": str(e)}


def generate_exercises(params):
    """Generate exercise list document"""
    prompt = _build_exercise_prompt(params)
    
    try:
        # Get content from LLM
        content = get_llm_response(prompt, params["llm_config"])
        
        # Create Word document
        docx_file = _create_exercise_docx(content, params)
        
        return {
            "success": True,
            "content": content,
            "file_data": docx_file
        }
        
    except Exception as e:
        return {"success": False, "error": str(e)}


def generate_powerpoint(params):
    """Generate PowerPoint presentation"""
    prompt = _build_powerpoint_prompt(params)
    
    try:
        # Get content from LLM
        content = get_llm_response(prompt, params["llm_config"])
        
        # Create PowerPoint file
        pptx_file = _create_powerpoint_pptx(content, params)
        
        return {
            "success": True,
            "content": content,
            "file_data": pptx_file
        }
        
    except Exception as e:
        return {"success": False, "error": str(e)}


def generate_summary(params):
    """Generate summary document"""
    prompt = _build_summary_prompt(params)
    
    try:
        # Get content from LLM
        content = get_llm_response(prompt, params["llm_config"])
        
        # Create Word document
        docx_file = _create_summary_docx(content, params)
        
        return {
            "success": True,
            "content": content,
            "file_data": docx_file
        }
        
    except Exception as e:
        return {"success": False, "error": str(e)}


def _build_exercise_prompt(params):
    """Build prompt for exercise generation"""
    prompt = f"""
Create a comprehensive exercise list for {params['subject']} at {params['grade_level']} level.

Topic: {params['topic']}
Number of questions: {params['num_questions']}
Difficulty: {params['difficulty']}
Question types: {', '.join(params['question_types'])}

Please structure the exercises as follows:
1. Start with a brief introduction to the topic
2. Organize questions by difficulty (if mixed difficulty is selected)
3. Include clear instructions for each section
4. For multiple choice questions, provide 4 options (A, B, C, D)
5. For problem-solving questions, show step-by-step solutions when appropriate
6. End with an answer key

Make sure the content is age-appropriate and educationally valuable.
Format the output in clear sections with proper headings.
"""
    return prompt


def _build_powerpoint_prompt(params):
    """Build prompt for PowerPoint generation"""
    prompt = f"""
Create a PowerPoint presentation outline for {params['subject']} at {params['grade_level']} level.

Topic: {params['topic']}
Number of slides: {params['num_slides']}
Include image placeholders: {params['include_images']}

Please structure the presentation as follows:
1. Title slide with topic and subject
2. Overview/agenda slide
3. Content slides covering the main points
4. Summary/conclusion slide

For each slide, provide:
- Slide title
- Key bullet points (3-5 per slide)
- Brief speaker notes
{"- Image suggestions (describe what type of image would be appropriate)" if params['include_images'] else ""}

Make the content engaging and appropriate for the target grade level.
Use clear, concise language and logical flow between slides.
"""
    return prompt


def _build_summary_prompt(params):
    """Build prompt for summary generation"""
    prompt = f"""
Create a {params['summary_length'].lower()} summary for {params['subject']} at {params['grade_level']} level.

Topic: {params['topic']}
Include examples: {params['include_examples']}

Please structure the summary as follows:
1. Introduction to the topic
2. Main concepts and key points
3. Important details and explanations
{"4. Relevant examples and applications" if params['include_examples'] else ""}
5. Conclusion with key takeaways

Make the summary:
- Educational and informative
- Appropriate for the target grade level
- Well-organized with clear sections
- {"Rich with practical examples" if params['include_examples'] else "Focused on core concepts"}

Use clear headings and bullet points where appropriate.
"""
    return prompt


def _create_exercise_docx(content, params):
    """Create Word document for exercises"""
    doc = Document()
    
    # Add title
    title = doc.add_heading(f"{params['subject']} - Exercise List", 0)
    title.alignment = 1  # Center alignment
    
    # Add subtitle
    doc.add_heading(f"Topic: {params['topic'][:50]}{'...' if len(params['topic']) > 50 else ''}", level=2)
    doc.add_paragraph(f"Grade Level: {params['grade_level']}")
    doc.add_paragraph(f"Difficulty: {params['difficulty']}")
    doc.add_paragraph("")
    
    # Add content
    _add_formatted_content_to_docx(doc, content)
    
    # Save to bytes
    file_buffer = io.BytesIO()
    doc.save(file_buffer)
    file_buffer.seek(0)
    return file_buffer.getvalue()


def _create_powerpoint_pptx(content, params):
    """Create PowerPoint presentation"""
    prs = Presentation()
    
    # Parse content and create slides
    slides_content = _parse_powerpoint_content(content)
    
    for slide_info in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_info.get('title', 'Slide Title')
        
        # Set content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.text = slide_info.get('content', 'Slide content')
    
    # Save to bytes
    file_buffer = io.BytesIO()
    prs.save(file_buffer)
    file_buffer.seek(0)
    return file_buffer.getvalue()


def _create_summary_docx(content, params):
    """Create Word document for summary"""
    doc = Document()
    
    # Add title
    title = doc.add_heading(f"{params['subject']} - Summary", 0)
    title.alignment = 1  # Center alignment
    
    # Add subtitle
    doc.add_heading(f"Topic: {params['topic'][:50]}{'...' if len(params['topic']) > 50 else ''}", level=2)
    doc.add_paragraph(f"Grade Level: {params['grade_level']}")
    doc.add_paragraph(f"Summary Type: {params['summary_length']}")
    doc.add_paragraph("")
    
    # Add content
    _add_formatted_content_to_docx(doc, content)
    
    # Save to bytes
    file_buffer = io.BytesIO()
    doc.save(file_buffer)
    file_buffer.seek(0)
    return file_buffer.getvalue()


def _add_formatted_content_to_docx(doc, content):
    """Add formatted content to Word document"""
    lines = content.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            doc.add_paragraph("")
            continue
            
        # Check if it's a heading (starts with #)
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        # Check if it's a bullet point
        elif line.startswith('- ') or line.startswith('* '):
            doc.add_paragraph(line[2:], style='List Bullet')
        # Check if it's a numbered list
        elif line[0].isdigit() and line[1:3] in ['. ', ') ']:
            doc.add_paragraph(line[3:], style='List Number')
        else:
            doc.add_paragraph(line)


def _parse_powerpoint_content(content):
    """Parse LLM content into structured slides"""
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Look for slide indicators
        if 'slide' in line.lower() and (':' in line or line.endswith(':')):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': line, 'content': ''}
        elif current_slide:
            current_slide['content'] += line + '\n'
    
    # Add the last slide
    if current_slide:
        slides.append(current_slide)
    
    # If no slides were parsed, create a single slide with all content
    if not slides:
        slides.append({
            'title': 'Content',
            'content': content
        })
    
    return slides
