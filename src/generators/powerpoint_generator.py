from llm_handlers.api_handler import get_llm_response
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import io
import re

def generate_powerpoint(params):
    """Generate PowerPoint presentation"""
    
    prompt = _build_powerpoint_prompt(params)
    
    try:
        # Get content from LLM
        content = get_llm_response(prompt, params["llm_config"])
        
        if not content or content.strip() == "":
            return {"success": False, "error": "LLM returned empty content"}
        
        # Create PowerPoint file
        pptx_file = _create_powerpoint_pptx(content, params)
        
        return {
            "success": True,
            "content": content,
            "pptx_file": pptx_file
        }
        
    except Exception as e:
        return {"success": False, "error": f"PowerPoint generation failed: {str(e)}"}

def _build_powerpoint_prompt(params):
    """Build prompt for PowerPoint generation"""
    
    prompt = f"""
    Create a {params['num_slides']}-slide PowerPoint presentation outline for {params['subject']} at {params['grade_level']} level.
    
    Topic: {params['topic']}
    Presentation style: {params['presentation_style']}
    Include images: {params['include_images']}
    
    IMPORTANT: Follow this EXACT format for each slide:
    
    SLIDE 1: [Slide Title Here]
    - First bullet point
    - Second bullet point
    - Third bullet point
    NOTES: Speaker notes for this slide
    {'IMAGE: Description of relevant image' if params['include_images'] else ''}
    
    SLIDE 2: [Next Slide Title]
    - First bullet point
    - Second bullet point
    - Third bullet point
    NOTES: Speaker notes for this slide
    {'IMAGE: Description of relevant image' if params['include_images'] else ''}
    
    Continue this pattern for all {params['num_slides']} slides.
    
    Guidelines:
    - Each slide should have 3-5 bullet points maximum
    - Keep bullet points concise and clear
    - Make content appropriate for {params['grade_level']} level
    - Use {params['presentation_style']} style
    - Include practical examples when possible
    - Ensure logical flow between slides
    
    Start your response with "SLIDE 1:" and follow the format exactly.
    """
    
    return prompt

def _create_powerpoint_pptx(content, params):
    """Create PowerPoint file from content"""
    
    try:
        prs = Presentation()
        
        # Parse content and create slides
        slides_data = _parse_powerpoint_content(content)
        
        if not slides_data:
            # Create a fallback slide if parsing fails
            slides_data = [{
                'title': f"{params['subject']} - {params.get('doc_type', 'Presentation')}",
                'bullets': ['Content could not be parsed properly', 'Please try generating again'],
                'notes': 'Generated content parsing failed',
                'image': ''
            }]
        
        for slide_data in slides_data:
            # Use title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title - check if title placeholder exists
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get('title', 'Slide Title')
            
            # Add content using a more robust approach
            bullets = slide_data.get('bullets', [])
            
            # Find the content placeholder
            content_shape = None
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:
                    content_shape = shape
                    break
            
            # If we found the content placeholder, add text to it
            if content_shape:
                try:
                    # Access text_frame safely with getattr
                    tf = getattr(content_shape, 'text_frame', None)
                    if tf:
                        tf.clear()
                        
                        if bullets:
                            # Set the first bullet point
                            tf.text = bullets[0]
                            
                            # Add additional bullet points
                            for bullet in bullets[1:]:
                                p = tf.add_paragraph()
                                p.text = bullet
                                p.level = 0
                        else:
                            tf.text = "No content available"
                    else:
                        raise AttributeError("No text_frame available")
                        
                except (AttributeError, Exception):
                    # If text_frame doesn't exist, add a text box instead
                    _add_text_box_to_slide(slide, bullets)
            else:
                # No content placeholder found, add a text box
                _add_text_box_to_slide(slide, bullets)
            
            # Add speaker notes if available
            notes_text = slide_data.get('notes', '')
            if notes_text:
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide and hasattr(notes_slide, 'notes_text_frame'):
                        text_frame = notes_slide.notes_text_frame
                        if text_frame:
                            text_frame.text = notes_text
                except Exception:
                    pass  # Skip notes if there's an issue
        
        # Save to bytes
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        return pptx_io.getvalue()
        
    except Exception as e:
        raise Exception(f"Failed to create PowerPoint file: {str(e)}")

def _add_text_box_to_slide(slide, bullets):
    """Add a text box with bullet points to a slide"""
    try:
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        if bullets:
            text_frame.text = bullets[0]
            for bullet in bullets[1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
        else:
            text_frame.text = "No content available"
    except Exception:
        pass  # Skip if text box creation fails

def _parse_powerpoint_content(content):
    """Parse LLM response into slide data with improved parsing"""
    
    slides = []
    current_slide = None
    
    lines = content.split('\n')
    
    for line in lines:
        line = line.strip()
        
        # Look for slide markers with various formats
        if re.match(r'^SLIDE\s*\d*:?', line, re.IGNORECASE) or line.upper().startswith('SLIDE'):
            # Save previous slide if exists
            if current_slide and (current_slide.get('title') or current_slide.get('bullets')):
                slides.append(current_slide)
            
            # Extract title from slide line
            title_match = re.search(r'SLIDE\s*\d*:?\s*(.+)', line, re.IGNORECASE)
            title = title_match.group(1).strip() if title_match else f"Slide {len(slides) + 1}"
            
            current_slide = {
                'title': title,
                'bullets': [],
                'notes': '',
                'image': ''
            }
            
        elif line.startswith('- ') or line.startswith('• ') or line.startswith('* '):
            # Handle bullet points
            if current_slide is not None:
                bullet_text = line[2:].strip()  # Remove bullet marker
                if bullet_text:
                    current_slide['bullets'].append(bullet_text)
                    
        elif line.upper().startswith('NOTES:'):
            # Handle speaker notes
            if current_slide is not None:
                notes_text = line[6:].strip()  # Remove "NOTES:" prefix
                current_slide['notes'] = notes_text
                
        elif line.upper().startswith('IMAGE:'):
            # Handle image suggestions
            if current_slide is not None:
                image_text = line[6:].strip()  # Remove "IMAGE:" prefix
                current_slide['image'] = image_text
                
        elif line and current_slide is not None and not current_slide.get('bullets'):
            # If we have text but no bullets yet, treat it as content
            # This handles cases where LLM doesn't use bullet format
            if not any(line.upper().startswith(prefix) for prefix in ['SLIDE', 'NOTES:', 'IMAGE:']):
                current_slide['bullets'].append(line)
    
    # Don't forget the last slide
    if current_slide and (current_slide.get('title') or current_slide.get('bullets')):
        slides.append(current_slide)
    
    # If no slides were parsed, create a fallback slide with the raw content
    if not slides:
        # Try to extract some meaningful content
        content_lines = [line.strip() for line in lines if line.strip() and not line.strip().startswith('#')]
        bullets = content_lines[:10] if content_lines else ["Unable to parse content"]
        
        slides.append({
            'title': 'Generated Content',
            'bullets': bullets,
            'notes': 'Content was parsed as fallback',
            'image': ''
        })
    
    return slides