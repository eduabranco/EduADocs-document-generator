#!/usr/bin/env python3
"""
Simplified Teacher Document Generator
A streamlined Streamlit application for generating educational documents using AI.
"""

import streamlit as st
import requests
import json
import os
from io import BytesIO
from docx import Document
from pptx import Presentation
from pptx.util import Inches


# Configuration
DEFAULT_MODELS = {
    "openai": ["gpt-4o-mini", "gpt-4o", "gpt-4-turbo", "gpt-3.5-turbo"],
    "ollama": ["llama2", "mistral", "codellama"],
    "huggingface": ["google/gemma-3-270m", "google/flan-t5-large"]
}


def get_api_key(service):
    """Get API key from environment or secrets"""
    try:
        return st.secrets.get(f"{service.upper()}_API_KEY", "")
    except:
        return os.getenv(f"{service.upper()}_API_KEY", "")


def check_ollama_connection(host="http://localhost:11434"):
    """Check if Ollama is running"""
    try:
        response = requests.get(f"{host}/api/tags", timeout=3)
        if response.status_code == 200:
            models = [m["name"] for m in response.json().get("models", [])]
            return True, models
        return False, []
    except:
        return False, []


def get_llm_response(prompt, config):
    """Unified function to get LLM response"""
    provider = config["provider"]
    
    if provider == "openai":
        return _call_openai(prompt, config)
    elif provider == "ollama":
        return _call_ollama(prompt, config)
    elif provider == "huggingface":
        return _call_huggingface(prompt, config)
    else:
        raise ValueError(f"Unsupported provider: {provider}")


def _call_openai(prompt, config):
    """Call OpenAI API"""
    headers = {
        "Authorization": f"Bearer {config['api_key']}",
        "Content-Type": "application/json"
    }
    
    data = {
        "model": config["model"],
        "messages": [{"role": "user", "content": prompt}],
        "temperature": config.get("temperature", 0.7),
        "max_tokens": 4000
    }
    
    response = requests.post(
        "https://api.openai.com/v1/chat/completions",
        headers=headers,
        json=data,
        timeout=60
    )
    
    if response.status_code != 200:
        raise Exception(f"OpenAI API error: {response.status_code}")
    
    result = response.json()
    return result["choices"][0]["message"]["content"]


def _call_ollama(prompt, config):
    """Call Ollama local API"""
    data = {
        "model": config["model"],
        "prompt": prompt,
        "stream": False,
        "options": {"temperature": config.get("temperature", 0.7)}
    }
    
    response = requests.post(
        f"{config['host']}/api/generate",
        json=data,
        timeout=300
    )
    
    if response.status_code != 200:
        raise Exception(f"Ollama error: {response.status_code}")
    
    return response.json().get("response", "No response generated")


def _call_huggingface(prompt, config):
    """Call Hugging Face API"""
    headers = {"Content-Type": "application/json"}
    if config.get("api_key"):
        headers["Authorization"] = f"Bearer {config['api_key']}"
    
    data = {
        "inputs": prompt,
        "parameters": {
            "temperature": config.get("temperature", 0.7),
            "max_new_tokens": 2000,
            "return_full_text": False
        }
    }
    
    response = requests.post(
        f"https://api-inference.huggingface.co/models/{config['model']}",
        headers=headers,
        json=data,
        timeout=60
    )
    
    if response.status_code != 200:
        raise Exception(f"Hugging Face error: {response.status_code}")
    
    result = response.json()
    if isinstance(result, list) and len(result) > 0:
        return result[0].get("generated_text", str(result[0]))
    return str(result)


def build_prompt(doc_type, params):
    """Build prompt based on document type and parameters"""
    base_prompt = f"""
Create a {doc_type.lower()} for {params['subject']} at {params['grade_level']} level.

Topic: {params['topic']}
"""
    
    if doc_type == "Exercise List":
        base_prompt += f"""
Number of questions: {params['num_questions']}
Difficulty: {params['difficulty']}
Question types: {', '.join(params['question_types'])}

Please create a well-structured exercise list with:
1. Clear instructions
2. Varied question types as requested
3. Appropriate difficulty level
4. Answer key at the end

Format the output in a clear, organized manner suitable for students.
"""
    
    elif doc_type == "PowerPoint Presentation":
        base_prompt += f"""
Number of slides: {params['num_slides']}
Style: {params['presentation_style']}
Include images: {params['include_images']}

Create a presentation outline with:
- Slide titles
- 3-5 bullet points per slide
- Speaker notes for each slide
{'- Image suggestions where appropriate' if params['include_images'] else ''}

Format as:
SLIDE 1: [Title]
- Bullet point 1
- Bullet point 2
NOTES: [Speaker notes]
{'IMAGE: [Image description]' if params['include_images'] else ''}
"""
    
    elif doc_type == "Summary":
        base_prompt += f"""
Length: {params['summary_length']}
Format: {params['format_style']}
Include examples: {params['include_examples']}

Create a comprehensive summary with:
1. Introduction to the topic
2. Main concepts and key points
3. {'Real-world examples' if params['include_examples'] else 'Theoretical explanations'}
4. Key takeaways
5. Further reading suggestions

Format according to the {params['format_style']} style requested.
"""
    
    return base_prompt


def create_docx_file(content, params):
    """Create Word document from content"""
    doc = Document()
    
    # Title
    title = doc.add_heading(f"{params['subject']} - {params.get('doc_type', 'Document')}", 0)
    title.alignment = 1  # Center
    
    # Subtitle
    doc.add_heading(f"Topic: {params['topic']}", level=2)
    doc.add_heading(f"Grade Level: {params['grade_level']}", level=3)
    
    # Content
    paragraphs = content.split('\n\n')
    for paragraph in paragraphs:
        if paragraph.strip():
            doc.add_paragraph(paragraph.strip())
    
    # Save to bytes
    doc_io = BytesIO()
    doc.save(doc_io)
    doc_io.seek(0)
    return doc_io.getvalue()


def create_pptx_file(content, params):
    """Create PowerPoint presentation from content"""
    prs = Presentation()
    
    # Parse content for slides
    slides_data = []
    current_slide = {}
    
    lines = content.split('\n')
    for line in lines:
        line = line.strip()
        if line.startswith('SLIDE'):
            if current_slide:
                slides_data.append(current_slide)
            current_slide = {
                'title': line.split(':', 1)[1].strip() if ':' in line else 'Slide',
                'bullets': [],
                'notes': ''
            }
        elif line.startswith('- ') and current_slide:
            current_slide['bullets'].append(line[2:])
        elif line.startswith('NOTES:') and current_slide:
            current_slide['notes'] = line[6:].strip()
    
    if current_slide:
        slides_data.append(current_slide)
    
    # Create slides
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Slide Title')
        
        # Content
        if slide_data.get('bullets'):
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.text = slide_data['bullets'][0] if slide_data['bullets'] else ""
            
            for bullet in slide_data['bullets'][1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
        
        # Notes
        if slide_data.get('notes'):
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = slide_data['notes']
    
    # Save to bytes
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io.getvalue()


def validate_inputs(subject, topic, llm_config):
    """Validate user inputs"""
    if not subject or not subject.strip():
        return False, "Subject is required"
    
    if not topic or not topic.strip():
        return False, "Topic description is required"
    
    if not llm_config:
        return False, "AI model configuration is required"
    
    provider = llm_config.get("provider")
    if provider == "openai" and not llm_config.get("api_key"):
        return False, "OpenAI API key is required"
    elif provider == "ollama" and not llm_config.get("connected", False):
        return False, "Ollama is not connected"
    
    return True, "Valid"


def configure_llm():
    """Configure LLM in sidebar"""
    st.sidebar.header("🤖 AI Model Selection")
    
    provider = st.sidebar.selectbox(
        "Provider",
        ["OpenAI API", "Ollama (Local)", "Hugging Face"]
    )
    
    config = {"provider": provider.split()[0].lower()}
    
    if provider == "OpenAI API":
        api_key = get_api_key("openai")
        if not api_key:
            api_key = st.sidebar.text_input("OpenAI API Key", type="password")
        else:
            st.sidebar.success("✅ API Key found")
        
        model = st.sidebar.selectbox("Model", DEFAULT_MODELS["openai"])
        temperature = st.sidebar.slider("Temperature", 0.0, 1.0, 0.7)
        
        config.update({
            "api_key": api_key,
            "model": model,
            "temperature": temperature
        })
    
    elif provider == "Ollama (Local)":
        host = st.sidebar.text_input("Host", "http://localhost:11434")
        connected, models = check_ollama_connection(host)
        
        if connected:
            st.sidebar.success("✅ Ollama connected")
            model = st.sidebar.selectbox("Model", models if models else DEFAULT_MODELS["ollama"])
        else:
            st.sidebar.error("❌ Ollama not connected")
            model = st.sidebar.text_input("Model", "llama2")
        
        temperature = st.sidebar.slider("Temperature", 0.0, 1.0, 0.7)
        
        config.update({
            "host": host,
            "model": model,
            "temperature": temperature,
            "connected": connected
        })
    
    elif provider == "Hugging Face":
        api_key = get_api_key("huggingface")
        if not api_key:
            api_key = st.sidebar.text_input("HF API Key (Optional)", type="password")
        
        model = st.sidebar.selectbox("Model", DEFAULT_MODELS["huggingface"])
        temperature = st.sidebar.slider("Temperature", 0.0, 1.0, 0.7)
        
        config.update({
            "api_key": api_key,
            "model": model,
            "temperature": temperature
        })
    
    return config


def main():
    """Main application"""
    st.set_page_config(
        page_title="Teacher Document Generator",
        page_icon="📚",
        layout="wide"
    )
    
    st.title("📚 Teacher Document Generator")
    st.markdown("Generate educational documents using AI")
    
    # Configure LLM
    llm_config = configure_llm()
    
    # Main form
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.header("📋 Document Settings")
        subject = st.text_input("Subject", placeholder="e.g., Mathematics, History")
        grade_level = st.selectbox(
            "Grade Level",
            ["Elementary (K-5)", "Middle School (6-8)", "High School (9-12)", "College", "Not specified"]
        )
        
        doc_type = st.selectbox(
            "Document Type",
            ["Exercise List", "PowerPoint Presentation", "Summary"]
        )
        
        topic = st.text_area(
            "Topic Description",
            height=150,
            placeholder="Describe the topic, learning objectives, or content..."
        )
        
        # Document-specific options
        params = {
            "subject": subject,
            "grade_level": grade_level,
            "topic": topic,
            "doc_type": doc_type,
            "llm_config": llm_config
        }
        
        if doc_type == "Exercise List":
            params["num_questions"] = st.slider("Number of Questions", 5, 50, 15)
            params["difficulty"] = st.selectbox("Difficulty", ["Easy", "Medium", "Hard", "Mixed"])
            params["question_types"] = st.multiselect(
                "Question Types",
                ["Multiple Choice", "True/False", "Short Answer", "Essay", "Problem Solving"],
                default=["Multiple Choice", "Short Answer"]
            )
        
        elif doc_type == "PowerPoint Presentation":
            params["num_slides"] = st.slider("Number of Slides", 5, 30, 12)
            params["include_images"] = st.checkbox("Include image suggestions", True)
            params["presentation_style"] = st.selectbox("Style", ["Educational", "Interactive", "Formal", "Creative"])
        
        elif doc_type == "Summary":
            params["summary_length"] = st.selectbox("Length", ["Brief (1-2 pages)", "Detailed (3-5 pages)", "Comprehensive (5+ pages)"])
            params["include_examples"] = st.checkbox("Include examples", True)
            params["format_style"] = st.selectbox("Format", ["Bullet Points", "Paragraphs", "Outline", "Q&A Format"])
    
    with col2:
        st.header("🚀 Generate")
        
        if st.button("Generate Document", type="primary", use_container_width=True):
            # Validate
            is_valid, message = validate_inputs(subject, topic, llm_config)
            
            if not is_valid:
                st.error(message)
                return
            
            try:
                with st.spinner("Generating document..."):
                    # Build prompt and get response
                    prompt = build_prompt(doc_type, params)
                    content = get_llm_response(prompt, llm_config)
                    
                    # Create files
                    docx_file = create_docx_file(content, params)
                    pptx_file = None
                    if doc_type == "PowerPoint Presentation":
                        pptx_file = create_pptx_file(content, params)
                    
                    st.success("✅ Document generated!")
                    
                    # Preview
                    with st.expander("📄 Preview", expanded=True):
                        st.markdown(content)
                    
                    # Downloads
                    st.header("💾 Downloads")
                    
                    st.download_button(
                        "📄 Download Word Document",
                        data=docx_file,
                        file_name=f"{subject}_{doc_type.replace(' ', '_')}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
                    
                    if pptx_file:
                        st.download_button(
                            "📊 Download PowerPoint",
                            data=pptx_file,
                            file_name=f"{subject}_presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
            
            except Exception as e:
                st.error(f"Error: {str(e)}")
        
        # Help
        with st.expander("❓ Help"):
            st.markdown("""
            **Quick Start:**
            1. Choose your AI provider
            2. Set document type and options
            3. Describe your topic
            4. Click Generate!
            
            **For Ollama:**
            - Run: `ollama serve`
            - Pull model: `ollama pull llama2`
            """)


if __name__ == "__main__":
    main()
