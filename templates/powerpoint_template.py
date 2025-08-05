from pptx import Presentation
from pptx.util import Inches

def create_powerpoint_slide(title, content):
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]  # Choosing a layout with title and content

    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

    return presentation

def save_powerpoint(presentation, filename):
    presentation.save(filename)

def generate_powerpoint(subject, exercises):
    presentation = create_powerpoint_slide(f"Exercises on {subject}", "\n".join(exercises))
    return presentation